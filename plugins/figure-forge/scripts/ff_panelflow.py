"""Panel-flow figures: staged pipeline diagrams in the style journals expect for
a proposed architecture or workflow.

Each column is a stage with a coloured header bar; each stage holds stacked
content boxes carrying a bold title and a grey detail line, optionally joined by
downward arrows, with arrows between the stages themselves. An optional banner
runs along the bottom for a status caveat.

Text is wrapped by *measuring* every candidate line with the real font, so a
line can never overflow its box - the containment guarantee the flowchart
renderer gives for single labels, extended to multi-line prose.
"""
from __future__ import annotations

import textwrap
import matplotlib.pyplot as plt
import matplotlib.patches as mp
from matplotlib.textpath import TextPath
from matplotlib.font_manager import FontProperties

import ff_style as S
from ff_typography import tx

MM = 1 / 25.4

# palette - muted, print-safe, and consistent with S.INK/S.SPINE
TEAL_DK = "#15525E"
TEAL    = "#2E7D8E"
TEAL_LT = "#BCD8DE"
PANEL   = "#EFF5F7"
GREY    = "#4F5B5E"
HI_FILL = "#DDEDF1"
OR_FILL = "#FDF6E7"
OR_EDGE = "#D9A441"
OR_INK  = "#8A6212"

FS_HDR, FS_TITLE, FS_BODY, FS_BANNER = 9.5, 8.2, 7.6, 8.0

STYLES = {
    "plain": ("white",   TEAL_LT, S.INK),
    "hi":    (HI_FILL,   TEAL,    TEAL_DK),
    "warn":  (OR_FILL,   OR_EDGE, OR_INK),
}


def _font(size, bold):
    return FontProperties(family=plt.rcParams["font.family"][0], size=size,
                          weight="bold" if bold else "normal")


def text_width_mm(text, size, bold=False):
    """True rendered width of `text` in mm."""
    if not text:
        return 0.0
    return TextPath((0, 0), text, prop=_font(size, bold)).get_extents().width / 72 / MM


def wrap(text, width_mm, size, bold=False):
    """Greedy wrap that measures each candidate line; never overflows width_mm."""
    out = []
    for para in str(text).split("\n"):
        words = para.split()
        if not words:
            out.append("")
            continue
        line = words[0]
        for w in words[1:]:
            cand = f"{line} {w}"
            if text_width_mm(cand, size, bold) <= width_mm:
                line = cand
            else:
                out.append(line)
                line = w
        out.append(line)
    return out


def _lh(size):
    return size * 1.30 / 72 / MM


def _geometry(spec):
    W = spec.get("width_mm", 170.0)
    cols = spec["columns"]
    gap = spec.get("col_gap", 4.0)
    pad_o = spec.get("outer_pad", 2.0)
    col_pad = spec.get("col_pad", 2.2)
    cw = (W - 2 * pad_o - (len(cols) - 1) * gap) / len(cols)
    return W, cols, gap, pad_o, col_pad, cw, cw - 2 * col_pad


def _box_h(b, inner):
    pad, h = 1.6, 1.6
    avail = inner - 3.2
    if b.get("title"):
        h += len(wrap(b["title"], avail, FS_TITLE, True)) * _lh(FS_TITLE) + 0.6
    if b.get("body"):
        h += len(wrap(b["body"], avail, FS_BODY)) * _lh(FS_BODY)
    return h + pad


def normalise(spec: dict) -> dict:
    """Apply the sign/punctuation rules to every string in the spec, ONCE.

    It has to happen before measurement, not at draw time. This renderer wraps
    and width-checks each line in millimetres and then draws it; converting the
    text afterwards would mean the containment QC measured one string and the
    figure drew another — an en dash is wider than a hyphen, so a line that
    passed the check could overflow its box in the file.
    """
    out = dict(spec)
    if out.get("title"):
        out["title"] = tx(out["title"])
    if out.get("banner"):
        out["banner"] = tx(out["banner"])
    cols = []
    for c in out.get("columns", out.get("cols", [])) or []:
        c = dict(c)
        if c.get("header"):
            c["header"] = tx(c["header"])
        c["boxes"] = [
            {**b,
             **({"title": tx(b["title"])} if b.get("title") else {}),
             **({"body": tx(b["body"])} if b.get("body") else {})}
            for b in c.get("boxes", [])
        ]
        cols.append(c)
    if "columns" in out:
        out["columns"] = cols
    elif "cols" in out:
        out["cols"] = cols
    return out


def check(spec):
    """Containment QC. Returns a list of violations; empty means clean."""
    spec = normalise(spec)
    W, cols, gap, pad_o, col_pad, cw, inner = _geometry(spec)
    avail = inner - 3.2
    bad = []
    for c in cols:
        for line in wrap(c.get("header", ""), cw - 4, FS_HDR, True):
            if text_width_mm(line, FS_HDR, True) > cw - 4 + 1e-6:
                bad.append({"rule": "R1", "where": "header", "text": line})
        for b in c["boxes"]:
            for key, size, bold in (("title", FS_TITLE, True), ("body", FS_BODY, False)):
                if not b.get(key):
                    continue
                for line in wrap(b[key], avail, size, bold):
                    w = text_width_mm(line, size, bold)
                    if w > avail + 1e-6:
                        bad.append({"rule": "R1", "where": key, "text": line,
                                    "width_mm": round(w, 2), "avail_mm": round(avail, 2)})
    return bad


def panelflow(spec, *, title=None):
    """Render a panel-flow figure. Returns (fig, violations)."""
    spec = normalise(spec)
    title = tx(title) if title else title
    W, cols, gap, pad_o, col_pad, cw, inner = _geometry(spec)
    hdr_h = spec.get("header_h", 9.0)
    box_gap = spec.get("box_gap", 2.2)
    arrow_h = 3.0
    banner = spec.get("banner")
    banner_h = spec.get("banner_h", 7.0) if banner else 0.0

    heights = []
    for c in cols:
        h = hdr_h + 2.0
        for b in c["boxes"]:
            h += _box_h(b, inner) + box_gap + (arrow_h if b.get("arrow_below") else 0)
        heights.append(h + 2.0)
    body_h = max(heights)
    H = pad_o + body_h + (banner_h + 3.0 if banner_h else 0) + pad_o
    if title:
        H += 6.0

    fig = plt.figure(figsize=(W * MM, H * MM))
    fig.patch.set_facecolor("white")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W)
    ax.set_ylim(0, H)
    ax.axis("off")

    top = H - pad_o
    if title:
        ax.text(pad_o, top - 2.0, title, ha="left", va="top",
                fontsize=FS_HDR, fontweight="bold", color=S.INK)
        top -= 6.0

    for ci, c in enumerate(cols):
        x0 = pad_o + ci * (cw + gap)
        ax.add_patch(mp.FancyBboxPatch(
            (x0, top - body_h), cw, body_h,
            boxstyle="round,pad=0,rounding_size=1.6", lw=0.7,
            ec=TEAL_LT, fc=PANEL, zorder=1))
        ax.add_patch(mp.FancyBboxPatch(
            (x0, top - hdr_h), cw, hdr_h,
            boxstyle="round,pad=0,rounding_size=1.6", lw=0,
            fc=c.get("header_color", TEAL_DK), zorder=2))
        hl = wrap(c.get("header", ""), cw - 4, FS_HDR, True)
        for k, line in enumerate(hl):
            ax.text(x0 + cw / 2,
                    top - hdr_h / 2 + (len(hl) - 1) * _lh(FS_HDR) / 2 - k * _lh(FS_HDR),
                    line, ha="center", va="center", color="white",
                    fontsize=FS_HDR, fontweight="bold", zorder=3)

        y = top - hdr_h - 2.0
        bx = x0 + col_pad
        for b in c["boxes"]:
            bh = _box_h(b, inner)
            fc, ec, tc = STYLES.get(b.get("style", "plain"), STYLES["plain"])
            ax.add_patch(mp.FancyBboxPatch(
                (bx, y - bh), inner, bh,
                boxstyle="round,pad=0,rounding_size=1.2", lw=0.7,
                ec=ec, fc=fc, zorder=3))
            ty = y - 1.6
            for key, size, bold, col in (("title", FS_TITLE, True, tc),
                                         ("body", FS_BODY, False, GREY)):
                if not b.get(key):
                    continue
                for line in wrap(b[key], inner - 3.2, size, bold):
                    ty -= _lh(size) * 0.80
                    ax.text(bx + 1.6, ty, line, ha="left", va="baseline",
                            color=col, fontsize=size, fontweight="bold" if bold else "normal",
                            zorder=4)
                    ty -= _lh(size) * 0.20
                if key == "title":
                    ty -= 0.6
            y -= bh
            if b.get("arrow_below"):
                ax.annotate("", xy=(bx + inner / 2, y - arrow_h + 0.4),
                            xytext=(bx + inner / 2, y - 0.4),
                            arrowprops=dict(arrowstyle="-|>", color=TEAL, lw=0.9,
                                            mutation_scale=8), zorder=4)
                y -= arrow_h
            y -= box_gap

        if ci < len(cols) - 1 and spec.get("col_arrows", True):
            ax.annotate("", xy=(x0 + cw + gap - 0.4, top - body_h / 2),
                        xytext=(x0 + cw + 0.4, top - body_h / 2),
                        arrowprops=dict(arrowstyle="-|>", color=TEAL, lw=1.3,
                                        mutation_scale=11), zorder=5)

    if banner:
        ax.add_patch(mp.FancyBboxPatch(
            (pad_o, pad_o), W - 2 * pad_o, banner_h,
            boxstyle="round,pad=0,rounding_size=1.4", lw=0.8,
            ec=OR_EDGE, fc=OR_FILL, zorder=3))
        ax.text(W / 2, pad_o + banner_h / 2, banner, ha="center", va="center",
                color=OR_INK, fontsize=FS_BANNER, zorder=4)

    fig._panelflow_texts = [t for t in ax.texts]
    return fig, check(spec)


def save_pptx_exact(fig, path, dpi=600):
    """PPTX at the figure's exact physical size, with a native, editable text box
    over every label.

    The stock exporter maps label positions into a tight-bbox crop; a panel-flow
    figure is built to a fixed physical width and is exported uncropped, so the
    mapping here is against the full canvas instead.
    """
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    w_in, h_in = fig.get_size_inches()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, facecolor="white")
    buf.seek(0)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(w_in), Inches(h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(buf, 0, 0, width=Inches(w_in), height=Inches(h_in))

    rnd = fig.canvas.get_renderer()
    fdpi = fig.dpi
    for t in getattr(fig, "_panelflow_texts", []):
        if not t.get_text().strip():
            continue
        bb = t.get_window_extent(rnd)          # display px, origin bottom-left
        x_in = bb.x0 / fdpi
        y_in = h_in - bb.y1 / fdpi             # flip to top-left origin
        tb = slide.shapes.add_textbox(Inches(max(0, x_in)), Inches(max(0, y_in)),
                                      Inches(bb.width / fdpi), Inches(bb.height / fdpi))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = {"left": 1, "center": 2, "right": 3}.get(t.get_ha(), 1)
        run = para.add_run()
        run.text = t.get_text()
        run.font.size = Pt(t.get_fontsize())
        run.font.bold = str(t.get_fontweight()) in ("bold", "600", "700", "800", "900")
        col = t.get_color()
        if isinstance(col, str) and col.startswith("#"):
            run.font.color.rgb = RGBColor.from_string(col.lstrip("#").upper())
    prs.save(str(path))
    return len(getattr(fig, "_panelflow_texts", []))
