"""Export a finished figure to every format a journal might ask for.

  svg   editable-text vector master (labels stay editable in Illustrator/Inkscape)
  png   raster at 600 dpi (default) - flatten for submission portals
  tiff  LZW-compressed, 600 dpi, dpi tag embedded - the classic Nature raster ask
  pdf   vector, embedded fonts
  pptx  one slide sized to the real figure, high-res image + optional native,
        editable PowerPoint text boxes over each label

All formats crop to the SAME tight bounding box so nothing (axis titles, panel
letters, outside labels) is ever clipped, and so the pptx slide, its image and
the editable text overlay stay in register.
"""
from __future__ import annotations

from pathlib import Path

import ff_editable as ED

PAD_IN = 0.03  # padding around the tight bbox, inches


def _tight_bbox(fig):
    """Padded tight bbox in inches (figure coords, origin bottom-left)."""
    fig.canvas.draw()
    bb = fig.get_tightbbox(fig.canvas.get_renderer())
    return bb.padded(PAD_IN)


def save_figure(fig, paths: dict, dpi=600, label_boxes=None, tight=True):
    """paths: {fmt: Path}. label_boxes: optional list for pptx editable text.

    tight=False exports the canvas exactly as sized, which is what a figure
    built to a fixed physical width (e.g. a 170 mm journal page) needs - a
    tight bbox would silently change that width.
    """
    bbox = _tight_bbox(fig) if tight else None
    written = {}
    for fmt, p in paths.items():
        p = Path(p)
        if fmt == "svg":
            fig.savefig(p, format="svg", bbox_inches=bbox)   # editable text
            # matplotlib names exactly one font and calls its groups `text_7`.
            # Neither survives being opened on another machine as anything a
            # human can edit, so the file is hardened before anyone sees it.
            ED.harden_svg(p)
        elif fmt == "pdf":
            fig.savefig(p, format="pdf", bbox_inches=bbox)
        elif fmt == "png":
            fig.savefig(p, format="png", dpi=dpi, bbox_inches=bbox)
        elif fmt in ("tif", "tiff"):
            _save_tiff(fig, p, dpi, bbox)
        elif fmt == "pptx":
            _save_pptx(fig, p, dpi, bbox, label_boxes)
        elif fmt == "eps":
            fig.savefig(p, format="eps", bbox_inches=bbox)
        else:
            raise ValueError(f"unknown format: {fmt}")
        written[fmt] = str(p)
    return written


def audit_outputs(written: dict) -> dict:
    """Check the files that exist, not the rcParams that were meant to make them.

    `svg.fonttype: none` can be set and the SVG can still come out with outlined
    text — a stale rcParam, a backend that ignored it, a figure element that
    forced a path. The raster proof looks identical either way, so the only
    honest check is to read the vector file back.
    """
    out = {}
    for fmt, path in written.items():
        try:
            if fmt == "svg":
                out["svg"] = ED.audit_svg(path)
            elif fmt == "pdf":
                out["pdf"] = ED.audit_pdf(path)
        except Exception as exc:                     # never fatal
            out[fmt] = {"error": str(exc)}
    return out


def _save_tiff(fig, path, dpi, bbox):
    import io
    from PIL import Image
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches=bbox)
    buf.seek(0)
    Image.open(buf).convert("RGB").save(
        path, format="TIFF", compression="tiff_lzw", dpi=(dpi, dpi))


def _save_pptx(fig, path, dpi, bbox, label_boxes):
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    w_in, h_in = bbox.width, bbox.height
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches=bbox)
    buf.seek(0)

    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(buf, 0, 0, width=Inches(w_in), height=Inches(h_in))

    for lb in (label_boxes or []):
        tb = slide.shapes.add_textbox(Inches(lb["x_in"]), Inches(lb["y_in"]),
                                      Inches(lb["w_in"]), Inches(lb["h_in"]))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = lb["text"]
        r.font.size = Pt(lb.get("fontsize", 7))
        c = lb.get("color", "#222222").lstrip("#")
        r.font.color.rgb = RGBColor.from_string(c.upper())
    prs.save(path)


def label_boxes_for_pptx(fig, verifier, dpi=600):
    """Editable-text boxes in inches, mapped into the tight-bbox crop frame so
    they line up with the embedded image in PPTX."""
    bbox = _tight_bbox(fig)
    rnd = fig.canvas.get_renderer()
    fdpi = fig.dpi                               # window extents are in fig.dpi px
    boxes = []
    for lab in verifier.labels:
        bb = lab.artist.get_window_extent(rnd)   # display px, origin bottom-left
        x_in = bb.x0 / fdpi - bbox.x0            # inches from crop left
        y_in = bbox.y1 - bb.y1 / fdpi           # inches from crop top
        boxes.append({
            "text": lab.artist.get_text(),
            "x_in": max(0, x_in), "y_in": max(0, y_in),
            "w_in": bb.width / fdpi, "h_in": bb.height / fdpi,
            "fontsize": lab.artist.get_fontsize(), "color": "#222222",
        })
    return boxes


# ------------------------------------------------- external SVG rasterisation -
def rasterize_svg(svg_path, out_path, dpi=600, fmt="png"):
    """Rasterise an arbitrary SVG (e.g. one we were asked to fix) via cairosvg."""
    import cairosvg
    svg_path, out_path = str(svg_path), str(out_path)
    if fmt == "png":
        cairosvg.svg2png(url=svg_path, write_to=out_path, dpi=dpi)
    elif fmt == "pdf":
        cairosvg.svg2pdf(url=svg_path, write_to=out_path, dpi=dpi)
    elif fmt in ("tif", "tiff"):
        import io
        from PIL import Image
        png = cairosvg.svg2png(url=svg_path, dpi=dpi)
        Image.open(io.BytesIO(png)).convert("RGB").save(
            out_path, format="TIFF", compression="tiff_lzw", dpi=(dpi, dpi))
    else:
        raise ValueError(f"cannot rasterise to {fmt}")
    return out_path
